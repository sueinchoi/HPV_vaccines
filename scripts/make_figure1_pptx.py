"""
Editable PowerPoint version of Figure 1 (cohort selection flow diagram).

Academic CONSORT-style layout — monochrome (black borders, white fill,
black text). The two branches diverge from the source population:
  * Cohort A: vaccine-status ascertainment → 1:1 PSM → exposure filter
  * Cohort B: cervical surgery first → post-surgery vaccine status →
              1:4 matching → exposure filter
Output: Data/Figure1_CohortSelection.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Layout (inches) — landscape slide 13.33 × 7.5 (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W, SLIDE_H = 13.33, 7.5
BOX_W = 4.8
BOX_H_S, BOX_H_M, BOX_H_L = 0.65, 0.85, 1.20
LEFT_X, RIGHT_X = 3.2, 9.0
CENTER_X = (LEFT_X + RIGHT_X) / 2

# Vertical row centres
ROW_SOURCE  = 0.55
ROW_HEADER  = 1.50
ROW_STEP1   = 2.45
ROW_STEP2   = 3.45
ROW_STEP3   = 4.45
ROW_STEP4   = 5.45      # Cohort B only (extra step)
ROW_FINAL_A = 5.50
ROW_FINAL_B = 6.55

FILL_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
FILL_LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
FILL_FINAL_GRAY  = RGBColor(0xE6, 0xE6, 0xE6)
EDGE_BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
RGB_TEXT         = RGBColor(0x00, 0x00, 0x00)
RGB_ARROW        = RGBColor(0x1A, 0x1A, 0x1A)

FONT_NAME = 'Arial'

prs = Presentation()
prs.slide_width  = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def add_box(cx, cy, w, h, lines, fill, edge, *, fs_first=13, fs_rest=12,
            bold_first=False):
    left = Inches(cx - w / 2); top = Inches(cy - h / 2)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = edge
    shape.line.width = Pt(1.0)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
    tf.margin_top  = Inches(0.05); tf.margin_bottom = Inches(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = line
        run.font.name = FONT_NAME
        run.font.size = Pt(fs_first if i == 0 else fs_rest)
        run.font.color.rgb = RGB_TEXT
        run.font.bold = (i == 0 and bold_first)
    return shape


def add_arrow(x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGB_ARROW
    line.line.width = Pt(1.0)
    from pptx.oxml.ns import qn
    ln = line.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')


# ---------- Source ----------
add_box(CENTER_X, ROW_SOURCE, 4.6, BOX_H_S,
        ['Source population',
         'N = 32,969'],
        FILL_WHITE, EDGE_BLACK, fs_first=14, fs_rest=13, bold_first=True)

# Branch arrows
add_arrow(CENTER_X, ROW_SOURCE + BOX_H_S/2, LEFT_X,  ROW_HEADER - BOX_H_S/2)
add_arrow(CENTER_X, ROW_SOURCE + BOX_H_S/2, RIGHT_X, ROW_HEADER - BOX_H_S/2)

# ---------- Cohort headers ----------
add_box(LEFT_X, ROW_HEADER, BOX_W, BOX_H_S,
        ['Cohort A \u2014 chronic-disease safety'],
        FILL_LIGHT_GRAY, EDGE_BLACK, fs_first=13, bold_first=True)
add_box(RIGHT_X, ROW_HEADER, BOX_W, BOX_H_S,
        ['Cohort B \u2014 post-surgical efficacy'],
        FILL_LIGHT_GRAY, EDGE_BLACK, fs_first=13, bold_first=True)

# ---------- Cohort A: vaccine status drives eligibility ----------
add_arrow(LEFT_X, ROW_HEADER + BOX_H_S/2, LEFT_X, ROW_STEP1 - BOX_H_M/2)
add_box(LEFT_X, ROW_STEP1, BOX_W, BOX_H_M,
        ['HPV-vaccine status ascertained',
         'Vaccinated 2,156   Unvaccinated 30,813'],
        FILL_WHITE, EDGE_BLACK, fs_first=12, fs_rest=11)

add_arrow(LEFT_X, ROW_STEP1 + BOX_H_M/2, LEFT_X, ROW_STEP2 - BOX_H_M/2)
add_box(LEFT_X, ROW_STEP2, BOX_W, BOX_H_M,
        ['1:1 propensity-score matching',
         '(age, BMI, BP, smoking, residence)'],
        FILL_WHITE, EDGE_BLACK, fs_first=12, fs_rest=11)

add_arrow(LEFT_X, ROW_STEP2 + BOX_H_M/2, LEFT_X, ROW_STEP3 - BOX_H_M/2)
add_box(LEFT_X, ROW_STEP3, BOX_W, BOX_H_M,
        ['Primary exposure filter',
         '\u22652 doses + 3-month landmark'],
        FILL_WHITE, EDGE_BLACK, fs_first=12, fs_rest=11)

add_arrow(LEFT_X, ROW_STEP3 + BOX_H_M/2, LEFT_X, ROW_FINAL_A - BOX_H_L/2)
add_box(LEFT_X, ROW_FINAL_A, BOX_W, BOX_H_L,
        ['Final analytic Cohort A',
         'n = 2,776',
         'Vaccinated 1,396  /  Unvaccinated 1,380'],
        FILL_FINAL_GRAY, EDGE_BLACK, fs_first=13, fs_rest=12, bold_first=True)

# ---------- Cohort B: surgery FIRST, then post-surgery vaccine status ----------
add_arrow(RIGHT_X, ROW_HEADER + BOX_H_S/2, RIGHT_X, ROW_STEP1 - BOX_H_M/2)
add_box(RIGHT_X, ROW_STEP1, BOX_W, BOX_H_M,
        ['Cervical surgery (conization or hysterectomy)',
         'n = 6,890'],
        FILL_WHITE, EDGE_BLACK, fs_first=12, fs_rest=11)

add_arrow(RIGHT_X, ROW_STEP1 + BOX_H_M/2, RIGHT_X, ROW_STEP2 - BOX_H_M/2)
add_box(RIGHT_X, ROW_STEP2, BOX_W, BOX_H_M,
        ['Post-surgery HPV-vaccine status',
         '(vaccinated after surgery vs never)'],
        FILL_WHITE, EDGE_BLACK, fs_first=12, fs_rest=11)

add_arrow(RIGHT_X, ROW_STEP2 + BOX_H_M/2, RIGHT_X, ROW_STEP3 - BOX_H_M/2)
add_box(RIGHT_X, ROW_STEP3, BOX_W, BOX_H_M,
        ['1:4 matching',
         '(surgery method/year/age, BMI)'],
        FILL_WHITE, EDGE_BLACK, fs_first=12, fs_rest=11)

add_arrow(RIGHT_X, ROW_STEP3 + BOX_H_M/2, RIGHT_X, ROW_STEP4 - BOX_H_M/2)
add_box(RIGHT_X, ROW_STEP4, BOX_W, BOX_H_M,
        ['Primary exposure filter',
         '\u22652 doses + 3-month landmark'],
        FILL_WHITE, EDGE_BLACK, fs_first=12, fs_rest=11)

add_arrow(RIGHT_X, ROW_STEP4 + BOX_H_M/2, RIGHT_X, ROW_FINAL_B - BOX_H_L/2)
add_box(RIGHT_X, ROW_FINAL_B, BOX_W, BOX_H_L,
        ['Final analytic Cohort B',
         'n = 912',
         'Vaccinated 203  /  Unvaccinated 709'],
        FILL_FINAL_GRAY, EDGE_BLACK, fs_first=13, fs_rest=12, bold_first=True)

out_path = 'Data/Figure1_CohortSelection.pptx'
prs.save(out_path)
print(f'Saved: {out_path}')
